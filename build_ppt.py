#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
《哥 · 歌 · 窑》 主旨解读 — PPT 生成脚本
输出: 哥窑_主旨解读.pptx (16:9)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import math as _math

# ================= 调色板 =================
INK          = RGBColor(0x0E, 0x0D, 0x0B)   # 墨黑
INK_2        = RGBColor(0x15, 0x13, 0x0F)
INK_3        = RGBColor(0x1D, 0x1A, 0x15)
PAPER        = RGBColor(0xF4, 0xEC, 0xDC)   # 瓷白
PAPER_2      = RGBColor(0xD8, 0xCD, 0xB6)
PAPER_MID    = RGBColor(0xB5, 0xAB, 0x96)
SMOKE        = RGBColor(0x5A, 0x67, 0x70)
KILN         = RGBColor(0xC0, 0x8A, 0x3E)   # 窑火金
KILN_2       = RGBColor(0xE3, 0xB2, 0x5A)
KILN_3       = RGBColor(0x8A, 0x5A, 0x1D)
PLUM         = RGBColor(0x7F, 0xAE, 0xA1)   # 梅子青
PLUM_2       = RGBColor(0x5D, 0x8E, 0x82)
EARTH        = RGBColor(0x9A, 0x3A, 0x2A)   # 紫金土
EARTH_2      = RGBColor(0xC4, 0x4A, 0x30)

GOLD_THIN    = RGBColor(0xC0, 0x8A, 0x3E)
PAPER_THIN   = RGBColor(0x3A, 0x35, 0x2C)

# ================= 字体 =================
F_CN = "思源宋体"     # Noto Serif SC / Source Han Serif
F_CN_BOLD = "思源宋体"
F_EN = "Cormorant Garamond"
F_MONO = "JetBrains Mono"
F_FALLBACK_CN = ["思源宋体", "Source Han Serif SC", "Noto Serif CJK SC", "宋体", "SimSun"]
F_FALLBACK_HK = ["思源黑体", "Source Han Sans SC", "Noto Sans CJK SC", "黑体", "SimHei"]

# ================= 工具 =================
def set_run_font(run, name=F_CN, size=14, bold=False, italic=False, color=PAPER):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    # 中文字体设置
    rPr = run._r.get_or_add_rPr()
    # 移除已有的 eastAsia 设置
    for ea in rPr.findall(qn('a:ea')):
        rPr.remove(ea)
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', name)

def add_textbox(slide, left, top, width, height, anchor='t'):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    if anchor == 't':
        tf.vertical_anchor = MSO_ANCHOR.TOP
    elif anchor == 'm':
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == 'b':
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    return tb, tf

def add_para(tf, text, font=F_CN, size=14, bold=False, italic=False,
             color=PAPER, align=PP_ALIGN.LEFT, space_before=0, space_after=0):
    if not tf.paragraphs[0].text and len(tf.paragraphs) == 1:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    # 兼容直接传 Pt 对象或数字
    try:
        p.space_before = space_before
        p.space_after = space_after
    except Exception:
        p.space_before = Pt(int(space_before))
        p.space_after = Pt(int(space_after))
    r = p.add_run()
    r.text = text
    set_run_font(r, name=font, size=size, bold=bold, italic=italic, color=color)
    return p, r

def add_rect(slide, left, top, width, height, fill=INK, line=None, line_w=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w) if line_w else Pt(0.5)
    shp.shadow.inherit = False
    return shp

def add_line(slide, x1, y1, x2, y2, color=KILN, weight=0.75):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = color
    cn.line.width = Pt(weight)
    return cn

def fill_background(slide, color=INK):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                 slide.part.package.presentation_part.presentation.slide_width,
                                 slide.part.package.presentation_part.presentation.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # 移到最底
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg

def add_text_rect(slide, left, top, width, height, fill=None, line=None, line_w=0.5):
    rect = add_rect(slide, left, top, width, height, fill=fill if fill else INK, line=line, line_w=line_w)
    return rect

# ================= 波形曲线（剧情曲线） =================
def draw_audio_wave(slide, x_start, x_end, baseline_y, peak_h,
                    n_samples=240, color=KILN_2, weight=1.2,
                    envelope_points=None, opacity_alpha=None):
    """
    绘制一条类似音频波形的平滑曲线，呼应「哥=歌」主题。
    envelope_points: [(t, amplitude_factor), ...]  用于在指定位置调节振幅
    """
    pts = []
    for i in range(n_samples):
        t = i / (n_samples - 1)
        x = x_start + t * (x_end - x_start)
        # 多频叠加 + 缓慢包络
        wave = (
            0.42 * _math.sin(t * _math.pi * 3.0) +
            0.28 * _math.sin(t * _math.pi * 7.0 + 1.1) +
            0.18 * _math.sin(t * _math.pi * 13.0 + 0.4) +
            0.12 * _math.sin(t * _math.pi * 19.0)
        )
        # 自定义包络
        env = 1.0
        if envelope_points:
            for j in range(len(envelope_points) - 1):
                t0, a0 = envelope_points[j]
                t1, a1 = envelope_points[j+1]
                if t0 <= t <= t1:
                    k = (t - t0) / (t1 - t0) if t1 > t0 else 0
                    env = a0 * (1 - k) + a1 * k
                    break
        y_off = wave * env * 0.5  # 归一
        y = baseline_y - y_off * peak_h
        pts.append((int(x), int(y)))

    # 画线段
    for i in range(len(pts) - 1):
        x1, y1 = pts[i]
        x2, y2 = pts[i+1]
        cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        cn.line.color.rgb = color
        cn.line.width = Pt(weight)

    return pts

# ================= 概念阶梯（金字塔层级） =================
def draw_concept_ladder(slide, x_start, y_top, layer_height, total_width,
                        layers, color_line=KILN, color_label=PAPER):
    """
    居中绘制一个由宽到窄的概念阶梯（自下而上）。
    layers: [(label, sub, width_ratio, fill), ...]  从底到顶
    """
    n = len(layers)
    cur_y = y_top
    for i, (label, sub, ratio, fill) in enumerate(layers):
        w = int(total_width * ratio)
        x = int(x_start + (total_width - w) / 2)
        h = layer_height
        # 阶梯条
        add_rect(slide, x, cur_y, w, h, fill=fill)
        # 顶边高亮
        add_rect(slide, x, cur_y, w, Emu(20000), fill=color_line)
        # 文案
        tb, tf = add_textbox(slide, x, cur_y, w, h, anchor='m')
        p, r = add_para(tf, label, font=F_CN, size=13, bold=True,
                        color=color_label, align=PP_ALIGN.CENTER)
        p2, r2 = add_para(tf, sub, font=F_EN, size=9, italic=True,
                          color=SMOKE, align=PP_ALIGN.CENTER, space_before=Pt(0))
        cur_y += h + Emu(40000)  # 层间小间隙
    return cur_y

# ================= 演示文稿 =================
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height

blank_layout = prs.slide_layouts[6]

# 通用页码与进度条绘制函数
def add_page_indicator(slide, idx, total=5):
    # 顶部进度条
    bar = add_rect(slide, 0, 0, int(SW * idx / total), Emu(25400), fill=KILN)
    # 页码
    tb, tf = add_textbox(slide, SW - Inches(1.6), Inches(0.25), Inches(1.3), Inches(0.3))
    p, r = add_para(tf, f"{idx:02d} / {total:02d}", font=F_MONO, size=10,
                    color=KILN_2, align=PP_ALIGN.RIGHT)

def add_corner_marks(slide):
    # 四角小金色记号
    mark_size = Inches(0.18)
    pad = Inches(0.25)
    positions = [
        (pad, pad, 0, 0),
        (SW - pad, pad, -1, 0),
        (pad, SH - pad, 0, -1),
        (SW - pad, SH - pad, -1, -1),
    ]
    for x, y, dx, dy in positions:
        # 横线
        h = add_rect(slide, x + dx*mark_size, y + dy*Emu(0), mark_size, Emu(15000), fill=KILN)
        # 竖线
        v = add_rect(slide, x + dx*Emu(0), y + dy*mark_size, Emu(15000), mark_size, fill=KILN)

# =====================================================
# Slide 1 — 封面
# =====================================================
s1 = prs.slides.add_slide(blank_layout)
fill_background(s1, INK)

# 窑火光晕（多层椭圆近似）
for r, alpha_factor, color in [
    (Inches(4.5), 0.45, KILN_2),
    (Inches(3.4), 0.30, EARTH_2),
    (Inches(2.4), 0.18, EARTH),
]:
    glow = s1.shapes.add_shape(MSO_SHAPE.OVAL,
                                int(SW/2 - r), int(SH*0.7 - r),
                                r*2, r*2)
    glow.fill.solid()
    glow.fill.fore_color.rgb = color
    glow.line.fill.background()
    # 透明度
    sp = glow.fill.fore_color._xFill
    solidFill = sp
    srgb = solidFill.find(qn('a:srgbClr'))
    if srgb is not None:
        alpha = etree.SubElement(srgb, qn('a:alpha'))
        alpha.set('val', str(int(255 * alpha_factor)))

# 装饰：左右竖线
add_rect(s1, Inches(0.5), Inches(0.5), Emu(8000), SH - Inches(1.0), fill=PAPER_THIN)
add_rect(s1, SW - Inches(0.5) - Emu(8000), Inches(0.5), Emu(8000), SH - Inches(1.0), fill=PAPER_THIN)

# 顶部：副标
tb, tf = add_textbox(s1, Inches(1.2), Inches(0.8), Inches(6), Inches(0.3))
add_para(tf, "主 旨 解 读 · 第 一 辑", font=F_CN, size=11, color=PAPER_MID)

# 顶部：英文副标
tb, tf = add_textbox(s1, SW - Inches(3.0), Inches(0.8), Inches(2.0), Inches(0.3), anchor='t')
add_para(tf, "GE · SONG · KILN", font=F_EN, size=12, italic=True, color=SMOKE, align=PP_ALIGN.RIGHT)

# 主标题"哥·歌·窑"
tb, tf = add_textbox(s1, Inches(1.0), Inches(1.7), Inches(11.3), Inches(3.0))
p, r = add_para(tf, "哥", font=F_CN, size=240, bold=True, color=PAPER, align=PP_ALIGN.LEFT)
p2, r2 = add_para(tf, "·  歌  ·  窑", font=F_CN, size=240, bold=True,
                  color=KILN_2, align=PP_ALIGN.LEFT, space_before=0)
# 实际上是三个字一行排版，重新做
tf.clear()
# 第一行
p1, r1 = add_para(tf, "哥", font=F_CN, size=220, bold=True, color=PAPER, align=PP_ALIGN.CENTER)
p1.space_before = Pt(0)
# 第二行
p2, r2 = add_para(tf, "·  歌  ·", font=F_CN, size=120, bold=True,
                  color=KILN_2, align=PP_ALIGN.CENTER, space_before=Pt(20))
# 第三行
p3, r3 = add_para(tf, "窑", font=F_CN, size=220, bold=True, color=PAPER, align=PP_ALIGN.CENTER,
                  space_before=Pt(20))

# 副标题英文
tb, tf = add_textbox(s1, Inches(1.0), Inches(6.0), Inches(11.3), Inches(0.4))
add_para(tf, "GE  ·  SONG  ·  KILN", font=F_EN, size=14, italic=True,
         color=SMOKE, align=PP_ALIGN.CENTER, space_before=Pt(0))

# 主旨引言
tb, tf = add_textbox(s1, Inches(2.0), Inches(6.5), Inches(9.3), Inches(0.5))
add_para(tf, "哥窑之「哥」，非长幼之序，乃「歌」之古字 —— 窑火如歌，兄弟如线。",
         font=F_CN, size=15, color=PAPER, align=PP_ALIGN.CENTER)

# 印章
seal = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, SW - Inches(1.6), Inches(1.3), Inches(1.0), Inches(1.0))
seal.fill.solid()
seal.fill.fore_color.rgb = EARTH
seal.line.color.rgb = PAPER
seal.line.width = Pt(1)
# 印章文字
tb, tf = add_textbox(s1, SW - Inches(1.6), Inches(1.3), Inches(1.0), Inches(1.0), anchor='m')
add_para(tf, "章", font=F_CN, size=44, bold=True, color=PAPER, align=PP_ALIGN.CENTER)

# 修正印记条
tb, tf = add_textbox(s1, Inches(1.0), Inches(7.0), Inches(11.3), Inches(0.3))
add_para(tf, "修 正 版  ·  哥 哥 = 张  寄   ·   弟 弟 = 张  志  元",
         font=F_CN, size=11, color=KILN_2, align=PP_ALIGN.CENTER)

add_page_indicator(s1, 1)


# =====================================================
# Slide 2 — 总纲
# =====================================================
s2 = prs.slides.add_slide(blank_layout)
fill_background(s2, INK)

# 头部
tb, tf = add_textbox(s2, Inches(0.7), Inches(0.45), Inches(2.0), Inches(0.4))
add_para(tf, "壹", font=F_CN, size=14, italic=True, color=KILN_2)

tb, tf = add_textbox(s2, Inches(1.5), Inches(0.4), Inches(8), Inches(0.7))
p, r = add_para(tf, "总   纲", font=F_CN, size=28, bold=True, color=PAPER)
p2, r2 = add_para(tf, "  /  THE  THESIS", font=F_EN, size=13, italic=True, color=SMOKE,
                  space_before=Pt(0))

# 副说明
tb, tf = add_textbox(s2, Inches(1.5), Inches(1.05), Inches(8), Inches(0.3))
add_para(tf, "核心主旨不变 · 人物标签互换", font=F_CN, size=12, color=PAPER_MID)

# 主旨条
add_rect(s2, Inches(0.7), Inches(1.65), Inches(12), Inches(0.85), fill=INK_2)
add_rect(s2, Inches(0.7), Inches(1.65), Emu(40000), Inches(0.85), fill=KILN)

tb, tf = add_textbox(s2, Inches(0.95), Inches(1.78), Inches(2.0), Inches(0.6), anchor='m')
add_para(tf, "THESIS", font=F_EN, size=12, italic=True, color=KILN_2)

tb, tf = add_textbox(s2, Inches(2.6), Inches(1.78), Inches(10), Inches(0.6), anchor='m')
add_para(tf, "《哥窑》是关于「传承」的视觉史诗。", font=F_CN, size=18, color=PAPER)

# 两张对比卡（紧凑）
card_y = Inches(2.7)
card_h = Inches(2.95)
card_w = Inches(5.4)

# 哥哥 卡
add_rect(s2, Inches(0.7), card_y, card_w, card_h, fill=INK_2, line=PAPER_THIN, line_w=0.75)
# 顶部细线
add_rect(s2, Inches(0.7), card_y, card_w, Emu(20000), fill=PLUM)

# 角色标签
tb, tf = add_textbox(s2, Inches(1.0), card_y + Inches(0.2), Inches(4), Inches(0.3))
add_para(tf, "持 灯 者  ·  守 护", font=F_CN, size=11, color=PLUM)

# 大字"哥哥"
tb, tf = add_textbox(s2, Inches(1.0), card_y + Inches(0.55), Inches(4), Inches(0.85))
add_para(tf, "哥  哥", font=F_CN, size=46, bold=True, color=PLUM)

# 名字
tb, tf = add_textbox(s2, Inches(1.0), card_y + Inches(1.4), Inches(4), Inches(0.35))
add_para(tf, "张   寄", font=F_CN, size=18, bold=True, color=PAPER)
# 分割
add_rect(s2, Inches(1.0), card_y + Inches(1.8), Inches(3.5), Emu(10000), fill=PAPER_THIN)

# 特质
traits_e = [
    ("理性", "规划的锚点"),
    ("承担", "战火中接过嘱托"),
    ("守护", "火候失控时拉住弟弟"),
]
for i, (k, v) in enumerate(traits_e):
    y = card_y + Inches(1.92 + i*0.24)
    tb, tf = add_textbox(s2, Inches(1.0), y, Inches(1.0), Inches(0.25))
    add_para(tf, f"▸ {k}", font=F_CN, size=11, color=PLUM_2)
    tb, tf = add_textbox(s2, Inches(2.0), y, Inches(3.0), Inches(0.25))
    add_para(tf, f"· {v}", font=F_CN, size=10, color=PAPER_MID)

# 引文
tb, tf = add_textbox(s2, Inches(1.0), card_y + Inches(2.6), Inches(4.2), Inches(0.25))
add_para(tf, "「等！现在开窑，全毁了！」", font=F_CN, size=10, italic=True, color=PAPER)
tb, tf = add_textbox(s2, Inches(1.0), card_y + Inches(2.78), Inches(4.2), Inches(0.2))
add_para(tf, "—— 镜号 53", font=F_MONO, size=8, color=SMOKE)

# 中央轴
ax_x = Inches(6.65)
add_rect(s2, ax_x, card_y + Inches(0.3), Emu(10000), card_h - Inches(0.6), fill=KILN)
# 节点
tb, tf = add_textbox(s2, ax_x - Inches(0.3), card_y + Inches(0.55), Inches(0.9), Inches(0.3), anchor='m')
add_para(tf, "持 守", font=F_CN, size=10, color=KILN_2, align=PP_ALIGN.CENTER)
tb, tf = add_textbox(s2, ax_x - Inches(0.3), card_y + Inches(2.4), Inches(0.9), Inches(0.3), anchor='m')
add_para(tf, "开 创", font=F_CN, size=10, color=KILN_2, align=PP_ALIGN.CENTER)
# 中心点
dot = s2.shapes.add_shape(MSO_SHAPE.OVAL, ax_x - Inches(0.06), card_y + Inches(1.45) - Inches(0.06),
                           Inches(0.12), Inches(0.12))
dot.fill.solid(); dot.fill.fore_color.rgb = KILN_2
dot.line.fill.background()

# 弟弟 卡
add_rect(s2, Inches(7.2), card_y, card_w, card_h, fill=INK_2, line=PAPER_THIN, line_w=0.75)
add_rect(s2, Inches(7.2), card_y, card_w, Emu(20000), fill=EARTH_2)

tb, tf = add_textbox(s2, Inches(7.5), card_y + Inches(0.2), Inches(4), Inches(0.3))
add_para(tf, "行 路 人  ·  探 索", font=F_CN, size=11, color=EARTH_2)

tb, tf = add_textbox(s2, Inches(7.5), card_y + Inches(0.55), Inches(4), Inches(0.85))
add_para(tf, "弟  弟", font=F_CN, size=46, bold=True, color=EARTH_2)

tb, tf = add_textbox(s2, Inches(7.5), card_y + Inches(1.4), Inches(4), Inches(0.35))
add_para(tf, "张   志   元", font=F_CN, size=18, bold=True, color=PAPER)
add_rect(s2, Inches(7.5), card_y + Inches(1.8), Inches(3.5), Emu(10000), fill=PAPER_THIN)

traits_y = [
    ("感性", "能听见「土在唱歌」"),
    ("天赋", "紫金土的发现者"),
    ("开创", "金丝铁线的引线人"),
]
for i, (k, v) in enumerate(traits_y):
    y = card_y + Inches(1.92 + i*0.24)
    tb, tf = add_textbox(s2, Inches(7.5), y, Inches(1.0), Inches(0.25))
    add_para(tf, f"▸ {k}", font=F_CN, size=11, color=EARTH_2)
    tb, tf = add_textbox(s2, Inches(8.5), y, Inches(3.0), Inches(0.25))
    add_para(tf, f"· {v}", font=F_CN, size=10, color=PAPER_MID)

tb, tf = add_textbox(s2, Inches(7.5), card_y + Inches(2.6), Inches(4.2), Inches(0.25))
add_para(tf, "「哥，这土……在唱歌。」", font=F_CN, size=10, italic=True, color=PAPER)
tb, tf = add_textbox(s2, Inches(7.5), card_y + Inches(2.78), Inches(4.2), Inches(0.2))
add_para(tf, "—— 镜号 16", font=F_MONO, size=8, color=SMOKE)

# 底部：剧情曲线（音频波形）— 呼应「哥 = 歌」主题
arc_label_y = Inches(5.7)
tb, tf = add_textbox(s2, Inches(0.7), arc_label_y, Inches(4.0), Inches(0.3))
add_para(tf, "三 幕  ·  情 绪 曲 线", font=F_CN, size=12, bold=True, color=KILN_2)
tb, tf = add_textbox(s2, Inches(0.7), arc_label_y + Inches(0.25), Inches(4.0), Inches(0.25))
add_para(tf, "THE  EMOTIONAL  WAVEFORM", font=F_EN, size=9, italic=True, color=SMOKE)

# 基线
baseline_y = Inches(6.15)
peak_h = Inches(0.55)
x_start = Inches(0.8)
x_end = Inches(12.5)
add_rect(s2, x_start, baseline_y, x_end - x_start, Emu(8000), fill=PAPER_THIN)

# 绘制音频波
# 包络: 序幕低 → 第一幕升 → 第二幕峰 → 第三幕先降后高 → 余韵平
envelope = [
    (0.00, 0.30),  # 序幕 - 离愁
    (0.05, 0.30),
    (0.15, 0.55),  # 探索
    (0.30, 0.75),  # 第一幕 紫金土
    (0.45, 0.70),  # 建窑
    (0.55, 0.90),  # 第二幕 贡品 - 峰 1
    (0.65, 0.95),
    (0.72, 0.55),  # 火候危机 - 谷
    (0.82, 0.50),
    (0.88, 1.00),  # 第三幕 开片·并肩 - 峰 2
    (0.95, 0.75),  # 余韵
    (1.00, 0.65),
]
draw_audio_wave(s2, x_start, x_end, baseline_y, peak_h,
                n_samples=300, color=KILN_2, weight=1.2,
                envelope_points=envelope)

# 关键节点标注
beats = [
    (0.05, "序幕",        "1-5",   "离愁",      PLUM,      Inches(-0.05)),
    (0.32, "第一幕",      "6-35",  "南下·紫金土", PAPER,    Inches(-0.05)),
    (0.58, "第二幕",      "36-47", "贡品之印", KILN_2,    Inches(-0.05)),
    (0.85, "第三幕",      "48-66", "开片·并肩", EARTH_2,   Inches(-0.05)),
]
for t, name, mirror, mood, col, yoff in beats:
    bx = Inches(0.8 + t * 11.7)
    # 节点圆点
    d = s2.shapes.add_shape(MSO_SHAPE.OVAL,
                             int(bx - Inches(0.07)),
                             int(baseline_y - Inches(0.07)),
                             Inches(0.14), Inches(0.14))
    d.fill.solid(); d.fill.fore_color.rgb = col
    d.line.color.rgb = INK
    d.line.width = Pt(1.5)
    # 标签在波的上方
    tb, tf = add_textbox(s2, int(bx - Inches(1.1)), int(baseline_y - Inches(0.65) + yoff),
                         Inches(2.2), Inches(0.6), anchor='b')
    p, r = add_para(tf, name, font=F_CN, size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
    p2, r2 = add_para(tf, mirror + "  ·  " + mood, font=F_CN, size=9, color=PAPER_MID,
                      align=PP_ALIGN.CENTER, space_before=Pt(0))

# 底部注脚
tb, tf = add_textbox(s2, Inches(0.7), Inches(7.05), Inches(12), Inches(0.3))
add_para(tf, "哥哥，是「歌」的定音鼓；弟弟，是「歌」的延长线 —— 整部短片，就是一首起伏绵长的歌。",
         font=F_CN, size=12, italic=True, color=PAPER_MID, align=PP_ALIGN.CENTER)
add_rect(s2, Inches(0.7), Inches(6.97), Inches(12), Emu(8000), fill=PAPER_THIN)

add_page_indicator(s2, 2)


# =====================================================
# Slide 3 — 证据链
# =====================================================
s3 = prs.slides.add_slide(blank_layout)
fill_background(s3, INK)

tb, tf = add_textbox(s3, Inches(0.7), Inches(0.45), Inches(2.0), Inches(0.4))
add_para(tf, "贰", font=F_CN, size=14, italic=True, color=KILN_2)

tb, tf = add_textbox(s3, Inches(1.5), Inches(0.4), Inches(10), Inches(0.7))
p, r = add_para(tf, "证   据   链", font=F_CN, size=28, bold=True, color=PAPER)
p2, r2 = add_para(tf, "  /  FROM THE SCRIPT", font=F_EN, size=13, italic=True, color=SMOKE)

tb, tf = add_textbox(s3, Inches(1.5), Inches(1.05), Inches(10), Inches(0.3))
add_para(tf, "五处镜号 · 两重性格的彼此成就", font=F_CN, size=12, color=PAPER_MID)

# 5 行证据
evidences = [
    ("哥哥·持守", PLUM,        "镜号 09",
     "「往南走，找能烧窑的土地，不能丢。手艺不能断。志元，看好你弟弟。」",
     "战乱中接过嘱托，以兄长之名承担全族手艺。"),
    ("弟弟·开创", EARTH_2,     "镜号 16",
     "「哥，这土……在唱歌。」",
     "紫金土的发现者 —— 天赋的感性回响。"),
    ("哥哥·持守", PLUM,        "镜号 17",
     "「那就这儿了。」",
     "以理性肯定弟弟的发现 —— 是节拍，亦是锚。"),
    ("弟弟·开创", EARTH_2,     "镜号 56",
     "「哥……成了。」",
     "金丝铁线开片，釉色如冰裂，释然之泪。"),
    ("弟弟→兄",   EARTH_2,     "镜号 45",
     "「哥，没有你这盏灯，我哪能坚持学业，这里都参着你的补贴呢。」",
     "两双手按在同一方印上 —— 二重唱成立的瞬间。"),
]

row_y0 = Inches(1.65)
row_h = Inches(0.95)
for i, (tag, tag_c, mirror, line, annot) in enumerate(evidences):
    y = row_y0 + i * row_h
    # 区分条
    add_rect(s3, Inches(0.7), y + row_h - Emu(8000), Inches(12), Emu(8000), fill=PAPER_THIN)
    # 标签
    add_rect(s3, Inches(0.7), y + Inches(0.2), Inches(1.4), Inches(0.45), fill=INK, line=tag_c, line_w=0.75)
    tb, tf = add_textbox(s3, Inches(0.7), y + Inches(0.2), Inches(1.4), Inches(0.45), anchor='m')
    add_para(tf, tag, font=F_CN, size=11, color=tag_c, align=PP_ALIGN.CENTER)
    # 镜号
    tb, tf = add_textbox(s3, Inches(2.3), y + Inches(0.2), Inches(1.2), Inches(0.4), anchor='m')
    add_para(tf, mirror, font=F_MONO, size=12, color=KILN_2)
    # 引文
    tb, tf = add_textbox(s3, Inches(3.6), y + Inches(0.1), Inches(9.0), Inches(0.4))
    p, r = add_para(tf, line, font=F_CN, size=15, italic=True, color=PAPER)
    # 注释
    tb, tf = add_textbox(s3, Inches(3.6), y + Inches(0.5), Inches(9.0), Inches(0.3))
    add_para(tf, annot, font=F_CN, size=11, color=SMOKE)

# 镜号 45 高亮强调
y_hi = row_y0 + 4 * row_h
hl = add_rect(s3, Inches(0.7), y_hi, Inches(12), row_h, fill=INK_3)
# 左侧强调条
add_rect(s3, Inches(0.7), y_hi, Emu(60000), row_h, fill=KILN)
# 重新覆盖该行
i = 4
tag, tag_c, mirror, line, annot = evidences[i]
y = y_hi
# 标签
add_rect(s3, Inches(0.9), y + Inches(0.2), Inches(1.4), Inches(0.45), fill=INK, line=tag_c, line_w=0.75)
tb, tf = add_textbox(s3, Inches(0.9), y + Inches(0.2), Inches(1.4), Inches(0.45), anchor='m')
add_para(tf, tag, font=F_CN, size=11, color=tag_c, align=PP_ALIGN.CENTER)
# 镜号
tb, tf = add_textbox(s3, Inches(2.4), y + Inches(0.2), Inches(1.2), Inches(0.4), anchor='m')
add_para(tf, mirror, font=F_MONO, size=12, color=KILN_2)
# 引文（高亮色）
tb, tf = add_textbox(s3, Inches(3.6), y + Inches(0.1), Inches(9.0), Inches(0.4))
p, r = add_para(tf, line, font=F_CN, size=15, italic=True, color=KILN_2)
# 注释
tb, tf = add_textbox(s3, Inches(3.6), y + Inches(0.5), Inches(9.0), Inches(0.3))
add_para(tf, "★ " + annot, font=F_CN, size=11, color=KILN_2)

add_page_indicator(s3, 3)


# =====================================================
# Slide 4 — 主旨升华
# =====================================================
s4 = prs.slides.add_slide(blank_layout)
fill_background(s4, INK)

tb, tf = add_textbox(s4, Inches(0.7), Inches(0.45), Inches(2.0), Inches(0.4))
add_para(tf, "叁", font=F_CN, size=14, italic=True, color=KILN_2)

tb, tf = add_textbox(s4, Inches(1.5), Inches(0.4), Inches(10), Inches(0.7))
p, r = add_para(tf, "主 旨 升 华", font=F_CN, size=28, bold=True, color=PAPER)
p2, r2 = add_para(tf, "  /  THE EPIPHANY", font=F_EN, size=13, italic=True, color=SMOKE)

tb, tf = add_textbox(s4, Inches(1.5), Inches(1.05), Inches(10), Inches(0.3))
add_para(tf, "哥  =  歌  =  兄", font=F_CN, size=12, color=PAPER_MID)

# 左侧：字源图（紧凑上移，留出底部阶梯区）
# "哥" 大字
tb, tf = add_textbox(s4, Inches(0.8), Inches(1.4), Inches(2.4), Inches(2.4), anchor='m')
add_para(tf, "哥", font=F_CN, size=160, bold=True, color=PAPER, align=PP_ALIGN.CENTER)
# 装饰横线
add_rect(s4, Inches(0.6), Inches(3.05), Inches(2.8), Emu(15000), fill=KILN)

# "兄" 印章
seal = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(1.9), Inches(1.1), Inches(1.1))
seal.fill.solid(); seal.fill.fore_color.rgb = EARTH
seal.line.color.rgb = PAPER
seal.line.width = Pt(1.5)
tb, tf = add_textbox(s4, Inches(3.5), Inches(1.9), Inches(1.1), Inches(1.1), anchor='m')
add_para(tf, "兄", font=F_CN, size=58, bold=True, color=PAPER, align=PP_ALIGN.CENTER)

# 等号
tb, tf = add_textbox(s4, Inches(4.7), Inches(2.2), Inches(0.8), Inches(0.6), anchor='m')
add_para(tf, "通", font=F_CN, size=22, bold=True, color=KILN_2, align=PP_ALIGN.CENTER)
# 箭头
arrow = s4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.9), Inches(2.4), Inches(0.7), Inches(0.18))
arrow.fill.solid(); arrow.fill.fore_color.rgb = KILN
arrow.line.fill.background()

# "歌" 古字
tb, tf = add_textbox(s4, Inches(5.7), Inches(1.4), Inches(2.4), Inches(2.4), anchor='m')
add_para(tf, "歌", font=F_CN, size=160, bold=True, color=KILN_2, align=PP_ALIGN.CENTER)
add_rect(s4, Inches(5.5), Inches(3.05), Inches(2.8), Emu(15000), fill=KILN)

# 字源说明
tb, tf = add_textbox(s4, Inches(1.0), Inches(3.2), Inches(6.5), Inches(0.3))
add_para(tf, "古 字 相 通  ·  字 源 之 证", font=F_CN, size=11, color=SMOKE, align=PP_ALIGN.CENTER)

# 右侧：三段升华文字（紧凑上移）
lines = [
    ("壹", "「哥」通「歌」 —— 窑火如歌，是声音，是节奏，是无数次的开与合。"),
    ("贰", "「哥」亦为「兄」 —— 兄弟如线，一根拉直，便成釉面上的金丝铁线。"),
    ("叁", "真正的传承，从来不是一个人的独唱 —— 而是两双手按在同一方印上、两个背影在时光中重叠的二重唱。"),
]
line_y0 = Inches(1.5)
for i, (num, text) in enumerate(lines):
    y = line_y0 + i * Inches(0.85)
    # 序号
    tb, tf = add_textbox(s4, Inches(8.4), y, Inches(0.7), Inches(0.8))
    add_para(tf, num, font=F_CN, size=28, bold=True, color=KILN)
    # 文字
    tb, tf = add_textbox(s4, Inches(9.2), y + Inches(0.05), Inches(3.8), Inches(0.85))
    p, r = add_para(tf, text, font=F_CN, size=11, color=PAPER, space_after=Pt(0))
    # 分隔
    if i < 2:
        add_rect(s4, Inches(8.4), y + Inches(0.75), Inches(4.6), Emu(6000), fill=PAPER_THIN)

# 底部：概念阶梯（升华主旨）
ladder_label_y = Inches(4.4)
tb, tf = add_textbox(s4, Inches(0.7), ladder_label_y, Inches(6.0), Inches(0.3))
add_para(tf, "概 念 阶 梯  ·  意 象 升 华", font=F_CN, size=12, bold=True, color=KILN_2)
tb, tf = add_textbox(s4, Inches(0.7), ladder_label_y + Inches(0.25), Inches(6.0), Inches(0.25))
add_para(tf, "THE  HIERARCHY  OF  INHERITANCE", font=F_EN, size=9, italic=True, color=SMOKE)

# 阶梯 4 层
ladder_layers = [
    # (label, sub, width_ratio, fill)
    ("传  ·  承",        "INHERITANCE  ·  哥窑之精神",  0.40, INK_3),
    ("哥 = 歌 = 兄",     "TRIUNITY  ·  三位一体",         0.62, INK_2),
    ("持 守  ×  开 创", "THE  DUET  ·  二重唱",         0.84, INK_2),
    ("兄 弟 二 人",      "THE  BROTHERS  ·  具象的人",  1.00, INK_3),
]
ladder_y = Inches(4.95)
ladder_h = Inches(0.4)
ladder_gap = Emu(40000)
ladder_total_w = Inches(10.0)
ladder_x_start = Inches((13.333 - 10.0) / 2)

# 绘制阶梯（自下而上：底层在最下，顶层在最上）
for i, (label, sub, ratio, fill) in enumerate(ladder_layers):
    w = int(ladder_total_w * ratio)
    x = int(ladder_x_start + (ladder_total_w - w) / 2)
    y = int(ladder_y + (len(ladder_layers) - 1 - i) * (ladder_h + ladder_gap))
    # 阶梯条
    add_rect(s4, x, y, w, ladder_h, fill=fill)
    # 顶边高亮
    add_rect(s4, x, y, w, Emu(20000), fill=KILN)
    # 文案
    tb, tf = add_textbox(s4, x, y, w, ladder_h, anchor='m')
    p, r = add_para(tf, label, font=F_CN, size=12, bold=True, color=PAPER, align=PP_ALIGN.CENTER)
    p2, r2 = add_para(tf, sub, font=F_EN, size=8, italic=True, color=KILN_2,
                      align=PP_ALIGN.CENTER, space_before=Pt(0))

# 阶梯上升箭头（左侧）
arrow_x = Inches(0.85)
arrow_y_start = Inches(7.0)
arrow_y_end = Inches(5.0)
add_rect(s4, arrow_x, arrow_y_start, Emu(20000), arrow_y_start - arrow_y_end, fill=KILN_2)
# 箭头三角
arrow_tip = s4.shapes.add_shape(MSO_SHAPE.UP_ARROW,
                                 int(arrow_x - Inches(0.06)),
                                 int(arrow_y_end - Inches(0.05)),
                                 Inches(0.17), Inches(0.2))
arrow_tip.fill.solid(); arrow_tip.fill.fore_color.rgb = KILN_2
arrow_tip.line.fill.background()
# 箭头旁注
tb, tf = add_textbox(s4, Inches(0.4), Inches(5.9), Inches(0.6), Inches(1.0), anchor='m')
add_para(tf, "升\n华", font=F_CN, size=10, bold=True, color=KILN_2, align=PP_ALIGN.CENTER)

# 阶梯右侧注
tb, tf = add_textbox(s4, Inches(11.5), Inches(5.8), Inches(1.6), Inches(1.2), anchor='m')
add_para(tf, "从\n具\n象\n到\n本\n质", font=F_CN, size=9, color=SMOKE, align=PP_ALIGN.CENTER)

# 印章融合进阶梯顶部"传承"层（中央印章感）
seal2 = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, SW - Inches(1.55), Inches(7.0), Inches(0.85), Inches(0.4))
seal2.fill.solid(); seal2.fill.fore_color.rgb = EARTH
seal2.line.color.rgb = PAPER
seal2.line.width = Pt(0.75)
tb, tf = add_textbox(s4, SW - Inches(1.55), Inches(7.0), Inches(0.85), Inches(0.4), anchor='m')
add_para(tf, "章  氏", font=F_CN, size=11, bold=True, color=PAPER, align=PP_ALIGN.CENTER)

add_page_indicator(s4, 4)


# =====================================================
# Slide 5 — 结语
# =====================================================
s5 = prs.slides.add_slide(blank_layout)
fill_background(s5, INK)

# 中心光晕
for r, alpha, color in [
    (Inches(4.0), 0.10, KILN_2),
    (Inches(3.0), 0.08, PLUM),
]:
    glow = s5.shapes.add_shape(MSO_SHAPE.OVAL,
                                int(SW/2 - r), int(SH*0.7 - r),
                                r*2, r*2)
    glow.fill.solid(); glow.fill.fore_color.rgb = color
    glow.line.fill.background()
    sp = glow.fill.fore_color._xFill
    srgb = sp.find(qn('a:srgbClr'))
    if srgb is not None:
        a = etree.SubElement(srgb, qn('a:alpha'))
        a.set('val', str(int(255 * alpha)))

# 顶部小标
tb, tf = add_textbox(s5, Inches(0), Inches(0.8), SW, Inches(0.4))
add_para(tf, "EPILOGUE   /   鸣 谢 与 指 正", font=F_EN, size=13, italic=True,
         color=KILN_2, align=PP_ALIGN.CENTER)

# 主引言
tb, tf = add_textbox(s5, Inches(1.5), Inches(1.6), SW - Inches(3.0), Inches(2.5), anchor='m')
add_para(tf, "「  窑火不灭，", font=F_CN, size=36, color=PAPER, align=PP_ALIGN.CENTER)
add_para(tf, "    兄弟如线。", font=F_CN, size=36, color=PAPER, align=PP_ALIGN.CENTER, space_before=Pt(8))
add_para(tf, "    一片瓷，", font=F_CN, size=36, color=KILN_2, align=PP_ALIGN.CENTER, space_before=Pt(8))
add_para(tf, "    便是一首歌的余韵。  」", font=F_CN, size=36, color=KILN_2, align=PP_ALIGN.CENTER, space_before=Pt(8))

# 元信息条
add_rect(s5, Inches(2.0), Inches(4.7), Inches(9.3), Emu(8000), fill=PAPER_THIN)
add_rect(s5, Inches(2.0), Inches(5.5), Inches(9.3), Emu(8000), fill=PAPER_THIN)

# 左：人物设定
tb, tf = add_textbox(s5, Inches(2.0), Inches(4.85), Inches(4.5), Inches(0.3))
add_para(tf, "人 物 设 定 · 修 正 后", font=F_EN, size=11, italic=True, color=SMOKE)
tb, tf = add_textbox(s5, Inches(2.0), Inches(5.15), Inches(4.5), Inches(0.3))
add_para(tf, "哥哥 = 张 寄   ·   弟弟 = 张 志 元",
         font=F_CN, size=14, color=PAPER)

# 右：核心主旨
tb, tf = add_textbox(s5, Inches(6.8), Inches(4.85), Inches(4.5), Inches(0.3))
add_para(tf, "核 心 主 旨", font=F_EN, size=11, italic=True, color=SMOKE)
tb, tf = add_textbox(s5, Inches(6.8), Inches(5.15), Inches(4.5), Inches(0.3))
add_para(tf, "哥 = 歌 = 兄   ·   持守 × 开创 的完美和弦",
         font=F_CN, size=14, color=PAPER)

# 鸣谢
tb, tf = add_textbox(s5, Inches(0), Inches(6.0), SW, Inches(0.4))
add_para(tf, "——", font=F_CN, size=14, color=KILN, align=PP_ALIGN.CENTER)

tb, tf = add_textbox(s5, Inches(0), Inches(6.2), SW, Inches(0.4))
add_para(tf, "非常感谢您的指正 —— 这是至关重要的细节。",
         font=F_CN, size=18, color=PAPER, align=PP_ALIGN.CENTER)

tb, tf = add_textbox(s5, Inches(0), Inches(6.65), SW, Inches(0.3))
add_para(tf, "文 末 谨 记  ·  修 订 印 记", font=F_CN, size=11, color=SMOKE, align=PP_ALIGN.CENTER)

tb, tf = add_textbox(s5, Inches(0), Inches(6.95), SW, Inches(0.4))
add_para(tf, "——", font=F_CN, size=14, color=KILN, align=PP_ALIGN.CENTER)

add_page_indicator(s5, 5)

# ================= 保存 =================
out = "/workspace/哥窑_主旨解读.pptx"
prs.save(out)
print(f"已生成：{out}")
print(f"幻灯片数：{len(prs.slides)}")
