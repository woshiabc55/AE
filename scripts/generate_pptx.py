#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成《哥窑》主旨解读 PPT 的可下载 .pptx 文件
基于 python-pptx，输出 24 张幻灯片，配色与 web 版本一致。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
import os

# ============ 颜色（与 web 端一致） ============
INK        = RGBColor(0x0c, 0x0a, 0x09)
INK_SOFT   = RGBColor(0x1a, 0x18, 0x14)
CELADON    = RGBColor(0x7e, 0xa8, 0x9a)
CELADON_L  = RGBColor(0xa8, 0xc4, 0xb8)
CELADON_D  = RGBColor(0x2f, 0x5b, 0x51)
KILN_GOLD  = RGBColor(0xd9, 0x9a, 0x52)
KILN_FLAME = RGBColor(0xf4, 0xa3, 0x5c)
KILN_EMBER = RGBColor(0xc0, 0x57, 0x3a)
PAPER      = RGBColor(0xf1, 0xeb, 0xe0)
PAPER_SOFT = RGBColor(0xe8, 0xe0, 0xd0)
GOLD_LINE  = RGBColor(0xb5, 0x8a, 0x4a)
ASH        = RGBColor(0x8a, 0x82, 0x75)

# 16:9 幻灯片
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]

# ============ 工具函数 ============
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.5)
    return s

def add_text(slide, x, y, w, h, text, *,
             font='SimSun', size=18, bold=False, color=PAPER,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # 中文 font
    rPr = run._r.get_or_add_rPr()
    eastAsian = rPr.find(qn('a:ea'))
    if eastAsian is None:
        from lxml import etree
        eastAsian = etree.SubElement(rPr, qn('a:ea'))
    eastAsian.set('typeface', 'SimSun')
    return tb

def add_multiline(slide, x, y, w, h, lines, *,
                  font='SimSun', color=PAPER, align=PP_ALIGN.LEFT, line_spacing=1.4):
    """lines: list of (text, size, bold, color)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, size, bold, c = item, 18, False, color
        elif len(item) == 4:
            text, size, bold, c = item
        else:
            text, size = item[0], item[1]
            bold = item[2] if len(item) > 2 else False
            c = item[3] if len(item) > 3 else color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = c
        from lxml import etree
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = etree.SubElement(rPr, qn('a:ea'))
        ea.set('typeface', 'SimSun')
    return tb

def add_eyebrow(slide, text, y=Inches(0.5)):
    """顶部小标签"""
    # 短线
    add_rect(slide, Inches(0.8), y + Inches(0.18), Inches(0.6), Pt(1), fill=KILN_GOLD)
    add_text(slide, Inches(1.5), y, Inches(10), Inches(0.4), text,
             font='Arial', size=11, color=KILN_GOLD, align=PP_ALIGN.LEFT)

def add_page_num(slide, num, total=24):
    txt = f'{num:02d} / {total:02d}'
    add_text(slide, Inches(11.3), Inches(0.3), Inches(1.5), Inches(0.3), txt,
             font='Arial', size=9, color=ASH, align=PP_ALIGN.RIGHT)

def add_progress_bar(slide, num, total=24):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Pt(2), fill=INK_SOFT)
    add_rect(slide, Inches(0), Inches(0),
             int(SLIDE_W * num / total), Pt(2), fill=KILN_GOLD)

# ============ 24 张幻灯片 ============

def slide_1():
    """封面"""
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, INK)
    add_progress_bar(s, 1)
    add_text(s, Inches(2), Inches(2.5), Inches(9), Inches(0.4),
             'A VISUAL EPIC OF HERITAGE · 视觉史诗',
             font='Arial', size=12, color=KILN_GOLD, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2), Inches(3), Inches(9), Inches(2.2),
             '哥 窑', font='SimSun', size=120, color=PAPER, align=PP_ALIGN.CENTER, bold=False)
    add_text(s, Inches(2), Inches(5.5), Inches(9), Inches(0.6),
             '窑火如歌 · 兄弟如线',
             font='SimSun', size=24, color=PAPER_SOFT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2), Inches(6.2), Inches(9), Inches(0.4),
             '主旨解读 · 主题演讲',
             font='Arial', size=10, color=ASH, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2), Inches(6.55), Inches(9), Inches(0.4),
             '24 SLIDES · 66 SHOTS · 3 ACTS',
             font='Arial', size=9, color=KILN_GOLD, align=PP_ALIGN.CENTER)
    add_page_num(s, 1)

def slide_2():
    """致辞"""
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, INK)
    add_progress_bar(s, 2)
    add_eyebrow(s, 'DEDICATION · 致 辞')
    add_multiline(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(2),
                  [('致 · 那只按住印的手', 56, False, KILN_GOLD)])
    add_multiline(s, Inches(0.8), Inches(3.5), Inches(11.5), Inches(1.5), [
        ('非常感谢您的指正。', 18, False, PAPER_SOFT),
        ('这是至关重要的细节。', 18, True, KILN_GOLD),
    ])
    add_text(s, Inches(0.8), Inches(5), Inches(11.5), Inches(0.5),
             '"哥哥是张寄，弟弟是张志元。"',
             font='SimSun', size=24, color=PAPER)
    add_text(s, Inches(0.8), Inches(5.7), Inches(11.5), Inches(0.4),
             '—— 主创手记',
             font='Arial', size=10, color=ASH)
    add_page_num(s, 2)

def slide_3():
    """目录"""
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, INK)
    add_progress_bar(s, 3)
    add_eyebrow(s, 'CONTENTS · 目 录')
    add_multiline(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.5),
                  [('本次演讲 · 五篇章', 48, False, PAPER)])
    items = [
        ('I',   '问句：哥窑的"哥"',       '04'),
        ('II',  '主旨：传承之歌',          '06'),
        ('III', '人物：两兄弟',            '08'),
        ('IV',  '第一幕：历史回溯',        '13'),
        ('V',   '第二幕：贡品之印',        '16'),
        ('VI',  '第三幕：窑火成瓷',        '19'),
        ('VII', '升华：持守与开创',        '22'),
        ('VIII','致敬：片尾',              '24'),
    ]
    y = Inches(3.2)
    for i, (num, title, page) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 5.7)
        yy = y + Inches(row * 0.7)
        add_text(s, x, yy, Inches(0.5), Inches(0.5), num,
                 font='Arial', size=11, color=KILN_GOLD, bold=True)
        add_text(s, x + Inches(0.5), yy, Inches(4), Inches(0.5), title,
                 font='SimSun', size=14, color=PAPER)
        add_text(s, x + Inches(4.5), yy, Inches(0.5), Inches(0.5), page,
                 font='Arial', size=10, color=ASH, align=PP_ALIGN.RIGHT)
        add_rect(s, x, yy + Inches(0.45), Inches(5), Pt(0.5), fill=KILN_GOLD)
    add_page_num(s, 3)

def slide_4():
    """序章·黑场"""
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, INK)
    add_progress_bar(s, 4)
    add_eyebrow(s, 'PROLOGUE · 序章 · 镜号 1–5')
    add_text(s, Inches(1), Inches(2), Inches(11), Inches(0.6),
             '黑场 · 特写 · 中景 · 大特写 · 黑场字幕',
             font='SimSun', size=20, color=ASH, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(3.5), Inches(11), Inches(1.5),
             '"哥窑的\'哥\', 是谁？"',
             font='SimSun', size=44, color=PAPER, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(5.2), Inches(11), Inches(0.4),
             '—— 镜号 5 · 黑场字幕',
             font='SimSun', size=12, color=ASH, align=PP_ALIGN.CENTER)
    add_page_num(s, 4)

def slide_5():
    """现代伏笔"""
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, INK)
    add_progress_bar(s, 5)
    add_eyebrow(s, 'THE FRAGMENT · 碎瓷之问')
    add_text(s, Inches(0.8), Inches(1.5), Inches(7), Inches(1.2),
             '缺了一块',
             font='SimSun', size=64, color=PAPER)
    add_multiline(s, Inches(0.8), Inches(3.5), Inches(7), Inches(3), [
        ('老匠人独坐桌前，看着瓷器，', 16, False, PAPER_SOFT),
        ('手滑险些掉落。', 16, True, KILN_GOLD),
        ('镜头另一端，角色在教育局', 16, False, PAPER_SOFT),
        ('申请教师资格和项目上报。', 16, False, PAPER_SOFT),
        ('', 8, False, PAPER),
        ('碎瓷底款——"章生"二字。', 18, True, KILN_GOLD),
    ])
    # 右侧大青瓷盘
    add_rect(s, Inches(9), Inches(2), Inches(3.5), Inches(3.5), fill=CELADON_D)
    add_text(s, Inches(9), Inches(3), Inches(3.5), Inches(1.5),
             '哥', font='SimSun', size=120, color=GOLD_LINE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(9), Inches(5.5), Inches(3.5), Inches(0.4),
             '碎瓷底款', font='SimSun', size=12, color=ASH, align=PP_ALIGN.CENTER)
    add_page_num(s, 5)

def slide_6():
    """主旨总纲"""
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, INK)
    add_progress_bar(s, 6)
    add_eyebrow(s, 'CORE THESIS · 主旨总纲')
    add_multiline(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(2.5), [
        ('一首关于', 56, False, PAPER),
        ('「传承」', 56, True, KILN_GOLD),
        ('的视觉史诗', 56, False, PAPER),
    ])
    add_rect(s, Inches(0.8), Inches(4.5), Pt(3), Inches(2), fill=KILN_GOLD)
    add_multiline(s, Inches(1.2), Inches(4.5), Inches(11), Inches(2.5), [
        ('"哥"通"歌"——', 28, False, PAPER),
        ('窑火如歌，兄弟如线。', 28, True, KILN_GOLD),
    ])
    add_page_num(s, 6)

def slide_7():
    """修正说明"""
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, INK)
    add_progress_bar(s, 7)
    add_eyebrow(s, 'REVISION · 关键修正')
    add_text(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.8),
             '人物关系的最终设定',
             font='SimSun', size=44, color=PAPER)
    # 表格
    rows = [
        ['角色', '正确设定', '核心特质', '视觉象征'],
        ['哥哥', '张寄', '守护者 · 理性 · 承担', '持灯者 · 锚点'],
        ['弟弟', '张志元', '探索者 · 感性 · 天赋', '行路人 · 延伸'],
    ]
    y = Inches(3.2)
    col_w = [Inches(1.5), Inches(2.5), Inches(3.5), Inches(4.0)]
    for ri, row in enumerate(rows):
        x = Inches(0.8)
        for ci, cell in enumerate(row):
            is_header = ri == 0
            color = KILN_GOLD if is_header else PAPER
            size = 11 if is_header else 14
            if ci == 0 and not is_header:
                color = CELADON_L if cell == '哥哥' else KILN_GOLD
            add_text(s, x, y, col_w[ci], Inches(0.5), cell,
                     font='SimSun' if is_header or ci < 2 else 'SimSun',
                     size=size, color=color)
            x += col_w[ci]
        if is_header:
            add_rect(s, Inches(0.8), y + Inches(0.45),
                     sum(col_w, Inches(0)), Pt(1), fill=KILN_GOLD)
        else:
            add_rect(s, Inches(0.8), y + Inches(0.45),
                     sum(col_w, Inches(0)), Pt(0.5), fill=ASH)
        y += Inches(0.8)
    add_page_num(s, 7)

def slide_8():
    """人物对照"""
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, INK)
    add_progress_bar(s, 8)
    add_eyebrow(s, 'CONTRAST · 人物对照')
    add_text(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.8),
             '哥哥 vs 弟弟',
             font='SimSun', size=44, color=PAPER)
    # 左卡片
    add_rect(s, Inches(0.8), Inches(3), Inches(5.5), Inches(3.8),
             fill=INK_SOFT, line=CELADON)
    add_text(s, Inches(1.1), Inches(3.2), Inches(5), Inches(0.6),
             '张寄', font='SimSun', size=36, color=PAPER)
    add_text(s, Inches(1.1), Inches(3.8), Inches(5), Inches(0.3),
             'THE ELDER · 哥哥', font='Arial', size=10, color=ASH)
    add_multiline(s, Inches(1.1), Inches(4.3), Inches(5), Inches(2), [
        ('· 持灯者 · 守护者 · 锚点', 14, False, PAPER_SOFT),
        ('· 理性 · 冷静 · 规划', 14, False, PAPER_SOFT),
        ('· 用树枝画出窑体轮廓', 14, False, PAPER_SOFT),
        ('· 火候失控时拉住弟弟', 14, False, PAPER_SOFT),
        ('· 把弟弟的手按在印上', 14, True, KILN_GOLD),
    ])
    # 右卡片
    add_rect(s, Inches(7), Inches(3), Inches(5.5), Inches(3.8),
             fill=INK_SOFT, line=KILN_GOLD)
    add_text(s, Inches(7.3), Inches(3.2), Inches(5), Inches(0.6),
             '张志元', font='SimSun', size=36, color=PAPER)
    add_text(s, Inches(7.3), Inches(3.8), Inches(5), Inches(0.3),
             'THE YOUNGER · 弟弟', font='Arial', size=10, color=ASH)
    add_multiline(s, Inches(7.3), Inches(4.3), Inches(5), Inches(2), [
        ('· 行路人 · 探索者 · 延伸', 14, False, PAPER_SOFT),
        ('· 感性 · 敏锐 · 天赋', 14, False, PAPER_SOFT),
        ('· 听见"土在唱歌"', 14, True, KILN_GOLD),
        ('· 听出"火太急了"', 14, False, PAPER_SOFT),
        ('· 开窑时泪与释然', 14, False, PAPER_SOFT),
    ])
    # 中间 VS
    add_text(s, Inches(6.2), Inches(4.5), Inches(1), Inches(0.5),
             'VS', font='Arial', size=24, color=ASH, align=PP_ALIGN.CENTER)
    add_page_num(s, 8)

def slide_simple_twocol(s, num, eyebrow, title_left, title_right,
                        left_items, right_items, left_color, right_color, note=None):
    add_progress_bar(s, num)
    add_eyebrow(s, eyebrow)
    add_text(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.7),
             title_left + ' / ' + title_right,
             font='SimSun', size=36, color=PAPER)
    # 左
    add_text(s, Inches(0.8), Inches(2.8), Inches(5.5), Inches(0.4),
             title_left, font='SimSun', size=20, color=left_color)
    add_multiline(s, Inches(0.8), Inches(3.3), Inches(5.5), Inches(3.5), [
        (it, 13, False, PAPER_SOFT) for it in left_items
    ])
    # 右
    add_text(s, Inches(7), Inches(2.8), Inches(5.5), Inches(0.4),
             title_right, font='SimSun', size=20, color=right_color)
    add_multiline(s, Inches(7), Inches(3.3), Inches(5.5), Inches(3.5), [
        (it, 13, False, PAPER_SOFT) for it in right_items
    ])
    if note:
        add_text(s, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.3),
                 note, font='SimSun', size=10, color=ASH)
    add_page_num(s, num)

def slide_9():
    """哥哥关键时刻"""
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, INK)
    slide_simple_twocol(
        s, 9, 'MOMENTS · 哥哥关键时刻', '张寄', '张志元',
        ['镜号 9  · 接过嘱托', '镜号 17 · 肯定 · "那就这儿了"',
         '镜号 23 · 默契 · 递来一块布', '镜号 53 · 拉住 · "等！"'],
        ['镜号 16 · 听土 · "在唱歌"', '镜号 29 · 听火 · "火着了"',
         '镜号 51 · 听危 · "太急了"', '镜号 57 · 听成 · "成了"'],
        CELADON_L, KILN_GOLD,
        note='哥哥 · 四个"守"的瞬间 / 弟弟 · 四个"听"的瞬间',
    )

def slide_10():
    """弟弟关键时刻"""
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, INK)
    slide_simple_twocol(
        s, 10, 'MOMENTS · 弟弟关键时刻', '张寄', '张志元',
        ['镜号 9  · 接过嘱托', '镜号 17 · 肯定 · "那就这儿了"',
         '镜号 23 · 默契 · 递来一块布', '镜号 53 · 拉住 · "等！"'],
        ['镜号 16 · 听土 · "在唱歌"', '镜号 29 · 听火 · "火着了"',
         '镜号 51 · 听危 · "太急了"', '镜号 57 · 听成 · "成了"'],
        CELADON_L, KILN_GOLD,
        note='持守者 × 探索者',
    )

def slide_11():
    """三幕结构"""
    s = prs.slides.add_slide(BLANK)
    set_slide_bg(s, INK)
    add_progress_bar(s, 11)
    add_eyebrow(s, 'STRUCTURE · 三幕结构')
    add_text(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.7),
             '三幕 · 一次"歌"的成形',
             font='SimSun', size=36, color=PAPER)
    # 三列
    cards = [
        ('I',   '历史回溯',  '镜号 6–35',  '寻土 · 起窑 · 第一把火\n未开片 · 兄弟奠基',          CELADON_L),
        ('II',  '贡品之印',  '镜号 36–47', '朝廷查验 · 天工之赐\n兄弟之印 · 全片高潮',        KILN_GOLD),
        ('III', '窑火成瓷',  '镜号 48–66', '火候失控 · 金丝铁线\n两个背影 · 闭环',             PAPER),
    ]
    x = Inches(0.8)
    for roman, title, range_, desc, c in cards:
        add_rect(s, x, Inches(2.8), Inches(3.7), Inches(3.5),
                 fill=INK_SOFT, line=KILN_GOLD)
        add_text(s, x + Inches(0.2), Inches(3), Inches(3), Inches(1.2),
                 roman, font='Arial', size=72, color=c)
        add_text(s, x + Inches(0.2), Inches(4.3), Inches(3.5), Inches(0.5),
                 title, font='SimSun', size=22, color=PAPER)
        add_text(s, x + Inches(0.2), Inches(4.8), Inches(3.5), Inches(0.3),
                 range_, font='Arial', size=10, color=ASH)
        add_text(s, x + Inches(0.2), Inches(5.2), Inches(3.5), Inches(1),
                 desc, font='SimSun', size=11, color=PAPER_SOFT)
        x += Inches(4.0)
    add_page_num(s, 11)

def slide_simple_act(s, num, eyebrow, act_text, title, body, quote, quote_cite):
    add_progress_bar(s, num)
    add_eyebrow(s, eyebrow)
    add_text(s, Inches(0.8), Inches(1.5), Inches(2), Inches(0.5),
             act_text, font='Arial', size=14, color=KILN_GOLD)
    add_text(s, Inches(0.8), Inches(2), Inches(11.5), Inches(0.7),
             title, font='SimSun', size=36, color=PAPER)
    add_multiline(s, Inches(0.8), Inches(3.2), Inches(11.5), Inches(2.5),
                  [(b, 14, False, PAPER_SOFT) for b in body])
    add_rect(s, Inches(0.8), Inches(5.2), Pt(3), Inches(1.5), fill=KILN_GOLD)
    add_text(s, Inches(1.2), Inches(5.2), Inches(11), Inches(1.5),
             quote, font='SimSun', size=20, color=PAPER)
    add_text(s, Inches(1.2), Inches(6.3), Inches(11), Inches(0.4),
             quote_cite, font='SimSun', size=10, color=ASH)
    add_page_num(s, num)

def slide_12():
    s = prs.slides.add_slide(BLANK); set_slide_bg(s, INK)
    slide_simple_act(s, 12, 'ACT ONE · 第一幕', 'ACT I',
                     '北方窑火 · 南下寻土',
                     ['战马、烽火、流民。',
                      '父亲的嘱托把两个包袱塞入兄弟手中——',
                      '一个是瓷土，一个是家谱。'],
                     '"往南走，找能烧窑的土地，不能丢。\n手艺不能断。志元，看好你弟弟。"',
                     '—— 镜号 9 · 父亲')

def slide_13():
    s = prs.slides.add_slide(BLANK); set_slide_bg(s, INK)
    slide_simple_act(s, 13, 'FIRST ACT · 寻土', 'ACT I',
                     '紫金土 · 听',
                     ['丽水群山，晨雾中土质泛红。',
                      '弟弟突然蹲下，抓起一把土，眼神发亮。',
                      '哥哥笑了——那笑容里有欣慰，有骄傲。'],
                     '"哥，这土……在唱歌。"\n"那就这儿了。"',
                     '—— 镜号 16/17 · 兄弟二人')

def slide_14():
    s = prs.slides.add_slide(BLANK); set_slide_bg(s, INK)
    slide_simple_act(s, 14, 'FIRST ACT · 起窑', 'ACT I',
                     '龙窑依山 · 兄弟夯土',
                     ['镜号 19 · 树枝画出窑体轮廓',
                      '镜号 21 · 并肩夯土，默契如一',
                      '镜号 23 · 哥哥默默递来一块布',
                      '镜号 25 · 向山民拱手承诺'],
                     '"龙窑要依山而建，火才能顺着山势往上走。"',
                     '—— 镜号 19 · 张寄')

def slide_15():
    s = prs.slides.add_slide(BLANK); set_slide_bg(s, INK)
    add_progress_bar(s, 15)
    add_eyebrow(s, 'FIRST ACT · 未开片')
    add_text(s, Inches(2), Inches(1.8), Inches(9), Inches(0.7),
             '第一窑 · 未开片', font='SimSun', size=48, color=PAPER, align=PP_ALIGN.CENTER)
    add_multiline(s, Inches(2), Inches(3), Inches(9), Inches(1.5), [
        ('烟散，第一件青瓷出窑——', 16, False, PAPER_SOFT),
        ('釉色如玉，却未开片。', 18, True, KILN_GOLD),
    ])
    add_text(s, Inches(2), Inches(5), Inches(9), Inches(0.5),
             '真正的杰作，永远不是一次成型的。',
             font='SimSun', size=20, color=PAPER, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2), Inches(5.6), Inches(9), Inches(0.4),
             '它是"守"与"走"的反复，',
             font='SimSun', size=14, color=PAPER_SOFT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2), Inches(6), Inches(9), Inches(0.4),
             '是失败与重来的合唱。',
             font='SimSun', size=14, color=PAPER_SOFT, align=PP_ALIGN.CENTER)
    add_page_num(s, 15)

def slide_16():
    s = prs.slides.add_slide(BLANK); set_slide_bg(s, INK)
    slide_simple_act(s, 16, 'ACT TWO · 第二幕', 'ACT II',
                     '朝廷查验 · 天工之赐',
                     ['几匹快马驰入山谷，官府旗帜飘扬。',
                      '一位官员下马——',
                      '朝廷有令，龙泉青瓷可贡于临安。'],
                     '"这纹路……如冰裂，却浑然天成？"',
                     '—— 镜号 40 · 官员')

def slide_17():
    s = prs.slides.add_slide(BLANK); set_slide_bg(s, INK)
    add_progress_bar(s, 17)
    add_eyebrow(s, 'SECOND ACT · 兄弟之印 ⭐')
    add_text(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.7),
             '两双手 · 一方印',
             font='SimSun', size=44, color=PAPER)
    add_rect(s, Inches(0.8), Inches(3), Pt(3), Inches(2.5), fill=KILN_GOLD)
    add_multiline(s, Inches(1.2), Inches(3), Inches(11), Inches(2.5), [
        ('张志元将印递给弟弟，', 22, False, PAPER),
        ('张寄却摇头，', 22, False, PAPER),
        ('把哥哥的手一起按在印上。', 22, True, KILN_GOLD),
    ])
    add_multiline(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5), [
        ('"哥，没有你这盏灯，我哪能坚持学业……"', 18, True, KILN_GOLD),
        ('—— 镜号 45 · 反向的致敬', 11, False, ASH),
    ])
    add_page_num(s, 17)

def slide_18():
    s = prs.slides.add_slide(BLANK); set_slide_bg(s, INK)
    add_progress_bar(s, 18)
    add_eyebrow(s, 'EVIDENCE · 贡品之印证据')
    add_text(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.7),
             '"无意之得 · 方为天工"',
             font='SimSun', size=36, color=PAPER)
    rows = [
        ['镜号 40', '官员眼中闪过惊讶，手指抚过釉面细纹', '细纹被发现 · 美的觉醒'],
        ['镜号 41', '张志元微颤答：是火候所致，非有意为之', '弟弟的谦逊 · 天赋不可强求'],
        ['镜号 42', '官员笑了，取出一方印', '朝廷认可 · 窑火登堂'],
        ['镜号 43', '印文"龙泉章氏"盖下', '姓氏之传 · 章氏之根'],
        ['镜号 45', '弟弟把哥哥的手按在印上', '⭐ 兄的谦让 · 二人的重'],
        ['镜号 47', '龙窑炊烟升起', '生机 · 传承的起点'],
    ]
    y = Inches(2.7)
    for r in rows:
        for ci, cell in enumerate(r):
            color = KILN_GOLD if ci == 0 else PAPER_SOFT
            add_text(s, Inches(0.8 + ci*4.1), y, Inches(4), Inches(0.4),
                     cell, font='SimSun', size=11, color=color)
        add_rect(s, Inches(0.8), y + Inches(0.4), Inches(12), Pt(0.5), fill=ASH)
        y += Inches(0.55)
    add_page_num(s, 18)

def slide_19():
    s = prs.slides.add_slide(BLANK); set_slide_bg(s, INK)
    slide_simple_act(s, 19, 'ACT THREE · 第三幕', 'ACT III',
                     '火候失控 · 金丝铁线',
                     ['张志元彻夜未眠，耳朵贴近窑壁——声音不对。',
                      '窑火突然转旺，窑孔中喷出金色火焰。',
                      '弟弟的疯狂遇上了哥哥的理性。'],
                     '"哥，火……太急了。"',
                     '—— 镜号 51 · 张志元')

def slide_20():
    s = prs.slides.add_slide(BLANK); set_slide_bg(s, INK)
    add_progress_bar(s, 20)
    add_eyebrow(s, 'THIRD ACT · 如一人 ⭐')
    add_text(s, Inches(0.8), Inches(1.5), Inches(7), Inches(0.7),
             '一人 · 一影',
             font='SimSun', size=48, color=PAPER)
    add_multiline(s, Inches(0.8), Inches(3), Inches(7), Inches(3), [
        ('兄弟二人并肩站在窑前，', 16, False, PAPER_SOFT),
        ('火光将他们的影子投在山壁上，', 16, True, KILN_GOLD),
        ('如一人。', 18, True, KILN_GOLD),
        ('', 8, False, PAPER),
        ('窑火声、心跳声、呼吸声——', 14, False, PAPER_SOFT),
        ('三个声音叠在一起，', 14, False, PAPER_SOFT),
        ('就是哥窑之"歌"。', 16, True, KILN_GOLD),
    ])
    # 右侧两个剪影
    cx = Inches(8.5)
    add_rect(s, cx, Inches(2.5), Inches(2), Inches(4),
             fill=INK_SOFT)
    # 火焰
    add_rect(s, cx + Inches(0.3), Inches(2.7), Inches(1.4), Inches(2.5),
             fill=KILN_FLAME)
    add_rect(s, cx + Inches(0.5), Inches(3), Inches(1), Inches(2),
             fill=KILN_GOLD)
    add_text(s, cx, Inches(6.5), Inches(2), Inches(0.3),
             '如一人', font='SimSun', size=11, color=ASH, align=PP_ALIGN.CENTER)
    add_page_num(s, 20)

def slide_21():
    s = prs.slides.add_slide(BLANK); set_slide_bg(s, INK)
    add_progress_bar(s, 21)
    add_eyebrow(s, 'THIRD ACT · 金丝铁线 ⭐')
    add_text(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.8),
             '金丝 · 铁线', font='SimSun', size=54, color=KILN_GOLD, align=PP_ALIGN.CENTER)
    # 仿开片方块
    add_rect(s, Inches(3), Inches(2.8), Inches(7.5), Inches(2.5), fill=CELADON)
    # 简易纹理
    from pptx.enum.shapes import MSO_CONNECTOR
    lines_coords = [
        (Inches(3.5), Inches(3.0), Inches(5), Inches(3.5)),
        (Inches(5), Inches(3.5), Inches(6.5), Inches(3.2)),
        (Inches(6.5), Inches(3.2), Inches(7.5), Inches(4.2)),
        (Inches(4), Inches(4), Inches(6), Inches(4.5)),
        (Inches(6), Inches(4.5), Inches(7.5), Inches(4.3)),
        (Inches(3.5), Inches(4.8), Inches(5.5), Inches(5.0)),
    ]
    for x1, y1, x2, y2 in lines_coords:
        c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        c.line.color.rgb = GOLD_LINE
        c.line.width = Pt(1.5)
    add_text(s, Inches(3), Inches(5.5), Inches(7.5), Inches(0.4),
             '金丝铁线 · 釉面开片',
             font='SimSun', size=11, color=ASH, align=PP_ALIGN.CENTER)
    add_multiline(s, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.8), [
        ('"哥……成了。"', 18, True, KILN_GOLD),
        ('—— 镜号 57 · 张志元 · 颤抖的手抚过瓷片，眼泪滑落', 10, False, ASH),
    ])
    add_page_num(s, 21)

def slide_22():
    s = prs.slides.add_slide(BLANK); set_slide_bg(s, INK)
    add_progress_bar(s, 22)
    add_eyebrow(s, 'RESOLUTION · 主旨升华')
    add_multiline(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(2.5), [
        ('"哥窑"的"哥"，是"歌"，也是"兄"。', 28, False, PAPER),
        ('哥哥张寄，就是那首"歌"的', 24, False, PAPER),
        ('定音鼓和节拍器。', 28, True, KILN_GOLD),
    ])
    # 双栏
    add_rect(s, Inches(0.8), Inches(4.8), Pt(3), Inches(1.5), fill=CELADON)
    add_multiline(s, Inches(1.2), Inches(4.8), Inches(5), Inches(1.5), [
        ('哥哥 · 张寄', 18, True, CELADON_L),
        ('持守 · 理性 · 守护', 12, False, PAPER_SOFT),
    ])
    add_rect(s, Inches(7), Inches(4.8), Pt(3), Inches(1.5), fill=KILN_GOLD)
    add_multiline(s, Inches(7.4), Inches(4.8), Inches(5), Inches(1.5), [
        ('弟弟 · 张志元', 18, True, KILN_GOLD),
        ('开创 · 感性 · 探索', 12, False, PAPER_SOFT),
    ])
    add_text(s, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.5),
             '哥窑之美，正在于持守与开创的完美和弦。',
             font='SimSun', size=20, color=PAPER, align=PP_ALIGN.CENTER)
    add_page_num(s, 22)

def slide_23():
    s = prs.slides.add_slide(BLANK); set_slide_bg(s, INK)
    add_progress_bar(s, 23)
    add_eyebrow(s, 'THE FINAL SENTENCE · 主题字幕')
    add_multiline(s, Inches(1), Inches(2), Inches(11), Inches(4), [
        ('哥窑之"哥"，', 32, False, PAPER),
        ('非长幼之序，', 32, False, PAPER),
        ('乃"歌"之古字——', 32, True, KILN_GOLD),
        ('窑火如歌，', 32, False, PAPER),
        ('兄弟如线。', 32, True, KILN_GOLD),
    ])
    add_text(s, Inches(1), Inches(6.5), Inches(11), Inches(0.4),
             '—— 镜号 64 · 黑场字幕 · 古琴最后一个音',
             font='SimSun', size=11, color=ASH, align=PP_ALIGN.CENTER)
    add_page_num(s, 23)

def slide_24():
    s = prs.slides.add_slide(BLANK); set_slide_bg(s, INK)
    add_progress_bar(s, 24)
    add_eyebrow(s, 'CREDITS · 致 谢')
    add_text(s, Inches(2), Inches(2.2), Inches(9), Inches(1.5),
             '哥 窑', font='SimSun', size=100, color=PAPER, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2), Inches(4.2), Inches(9), Inches(0.5),
             '窑火不灭 · 兄弟如线',
             font='SimSun', size=22, color=PAPER_SOFT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2), Inches(5.5), Inches(9), Inches(0.4),
             '24 SLIDES · 66 SHOTS · 3 ACTS',
             font='Arial', size=10, color=KILN_GOLD, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2), Inches(6.2), Inches(9), Inches(0.5),
             '— 谢 谢 —',
             font='SimSun', size=24, color=KILN_GOLD, align=PP_ALIGN.CENTER)
    add_page_num(s, 24)


# 全部生成
for fn in [slide_1, slide_2, slide_3, slide_4, slide_5, slide_6, slide_7, slide_8,
           slide_9, slide_10, slide_11, slide_12, slide_13, slide_14, slide_15,
           slide_16, slide_17, slide_18, slide_19, slide_20, slide_21, slide_22,
           slide_23, slide_24]:
    fn()

out = '/workspace/downloads/哥窑_主旨解读.pptx'
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print(f'已生成 {out} ({os.path.getsize(out)} bytes)')
